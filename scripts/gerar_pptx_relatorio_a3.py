"""
Gera o modelo de Relatório A3 (PPTX) com a identidade visual do Encontro GEI.

Saída: Modelos/RELATÓRIO A3 Encontro de Governança, Estratégia e Inovação_v2.pptx

Slide único 14"×10,5". Estrutura DMAIC em 5 caixas com headers nas cores da
marca (amarelo · verde · laranja · azul · navy), badges numerados 01–05 e
setas de fluxo indicando o ciclo. Conteúdo-orientação preservado verbatim
do template original.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import gerar_pptx_apresentacao as base

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
OUT = ROOT / "Modelos" / "RELATÓRIO A3 Encontro de Governança, Estratégia e Inovação_v2.pptx"

# Paleta — herda do template do vídeo + LARANJA do logo (Estratégia)
NAVY = base.NAVY
NAVY_2 = base.NAVY_2
YELLOW = base.YELLOW_BAR
GREEN = base.GREEN_BAR
BLUE = base.BLUE_BAR
ORANGE = RGBColor(0xF2, 0x91, 0x3D)  # tom da palavra "Estratégia" no logo
MUTED = base.MUTED
WHITE = base.WHITE
BG = base.BG
BG_SOFT = base.BG_SOFT
BG_BORDER = base.BG_BORDER
AUTHORS_BG = RGBColor(0xEE, 0xF1, 0xF8)  # navy diluído para a faixa de autores

FONT_HEAD = base.FONT_HEAD
FONT_BODY = base.FONT_BODY


def add_text(slide, x, y, w, h, text, *, size=11, bold=False, italic=False,
             color=NAVY, font=FONT_BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_color_stripe(slide, slide_w, y, *, height=Emu(76200)):
    """Faixa decorativa com 4 cores da marca: amarelo, verde, laranja, azul."""
    seg_w = slide_w // 4
    for i, c in enumerate((YELLOW, GREEN, ORANGE, BLUE)):
        seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, seg_w * i, y, seg_w, height)
        seg.line.fill.background()
        seg.fill.solid()
        seg.fill.fore_color.rgb = c


def add_image_safe(slide, path, x, y, *, w=None, h=None):
    if path.exists():
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    return None


def fill_rect(slide, x, y, w, h, color, *, line_color=None, line_width=None,
              shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_width is not None:
            sh.line.width = line_width
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def dmaic_box(slide, x, y, w, h, *, numero, etapa, header_color,
              body_paragraphs, header_h=Inches(0.55), is_dark_header=False):
    """Caixa DMAIC: header colorido + badge numerado + corpo branco.

    is_dark_header=True para barras escuras (NAVY): texto branco e badge
    branco com número navy. Default: barras claras → texto navy e badge
    navy com número branco (alto contraste com o pastel da marca).
    """
    # corpo (atrás)
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    body.adjustments[0] = 0.04
    body.line.color.rgb = BG_BORDER
    body.line.width = Pt(0.75)
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE

    # header (faixa colorida)
    hdr = fill_rect(slide, x, y, w, header_h, header_color)
    hdr.line.fill.background()

    if is_dark_header:
        title_color = WHITE
        badge_fill = WHITE
        badge_num_color = NAVY
        badge_line = WHITE
    else:
        title_color = NAVY
        badge_fill = NAVY
        badge_num_color = WHITE
        badge_line = NAVY

    # badge circular numerado
    badge_d = Inches(0.40)
    badge_x = x + Inches(0.16)
    badge_y = y + (header_h - badge_d) // 2
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y, badge_d, badge_d)
    badge.line.color.rgb = badge_line
    badge.line.width = Pt(1.5)
    badge.fill.solid()
    badge.fill.fore_color.rgb = badge_fill
    btf = badge.text_frame
    btf.margin_left = btf.margin_right = Emu(0)
    btf.margin_top = btf.margin_bottom = Emu(0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = f"{numero:02d}"
    br.font.name = FONT_HEAD
    br.font.size = Pt(13)
    br.font.bold = True
    br.font.color.rgb = badge_num_color

    # nome da etapa
    label_x = badge_x + badge_d + Inches(0.14)
    label = slide.shapes.add_textbox(
        label_x, y, w - (label_x - x) - Inches(0.16), header_h,
    )
    ltf = label.text_frame
    ltf.margin_left = ltf.margin_right = Emu(0)
    ltf.margin_top = ltf.margin_bottom = Emu(0)
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run()
    lr.text = etapa.upper()
    lr.font.name = FONT_HEAD
    lr.font.size = Pt(16)
    lr.font.bold = True
    lr.font.color.rgb = title_color

    # corpo de texto
    body_y = y + header_h + Inches(0.10)
    body_h = h - header_h - Inches(0.18)
    tb = slide.shapes.add_textbox(x + Inches(0.22), body_y, w - Inches(0.44), body_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, (text, opts) in enumerate(body_paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = opts.get("align", PP_ALIGN.LEFT)
        para.space_after = Pt(opts.get("space_after", 4))
        run = para.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(opts.get("size", 11))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", NAVY_2)
    return body


def cycle_arrow(slide, x, y, w, h, *, direction="down", color=None):
    color = color or NAVY_2
    shape_map = {
        "down": MSO_SHAPE.DOWN_ARROW,
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
    }
    arr = slide.shapes.add_shape(shape_map[direction], x, y, w, h)
    arr.line.fill.background()
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    return arr


def realizadoras_strip(slide, slide_w, y_top):
    add_text(
        slide, Inches(0.5), y_top, slide_w - Inches(1.0), Inches(0.22),
        "REALIZAÇÃO", size=8, bold=True, color=NAVY_2,
        font=FONT_HEAD, align=PP_ALIGN.CENTER,
    )

    logo_h = Inches(0.42)
    y_logos = y_top + Inches(0.26)
    spacing = Inches(0.5)
    slot_w = Inches(1.5)

    total_w = slot_w * 4 + spacing * 3
    start_x = (slide_w - total_w) // 2

    logo_uff = ASSETS / "logo-uff.png"
    if logo_uff.exists():
        add_image_safe(slide, logo_uff, start_x, y_logos, h=logo_h)
    else:
        base.add_uff_badge(slide, start_x, y_logos, slot_w, logo_h)

    add_image_safe(slide, ASSETS / "logos" / "abar.png",
                   start_x + (slot_w + spacing) * 1, y_logos, h=logo_h)
    add_image_safe(slide, ASSETS / "logos" / "labdge.png",
                   start_x + (slot_w + spacing) * 2, y_logos, h=logo_h)
    add_image_safe(slide, ASSETS / "logos" / "gigs-unicamp.jpg",
                   start_x + (slot_w + spacing) * 3, y_logos, h=logo_h)


def build():
    prs = Presentation()
    prs.slide_width = Inches(14)
    prs.slide_height = Inches(10.5)
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # fundo branco
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # faixa decorativa topo
    add_color_stripe(slide, SLIDE_W, 0, height=Emu(76200))

    # ===== Cabeçalho =====
    header_y = Inches(0.18)
    header_h = Inches(0.95)

    logo_evento = ASSETS / "logo-completo.png"
    if logo_evento.exists():
        add_image_safe(slide, logo_evento, Inches(0.45), header_y + Inches(0.05),
                       h=Inches(0.85))

    # título (verbatim do original)
    add_text(
        slide, Inches(2.4), header_y, Inches(8.6), header_h,
        "TÍTULO, Arial, 16, centralizado",
        size=20, bold=True, color=NAVY, font=FONT_HEAD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # marcador "RELATÓRIO A3 · CICLO DMAIC" à direita
    add_text(
        slide, SLIDE_W - Inches(2.6), header_y + Inches(0.10),
        Inches(2.2), Inches(0.32),
        "RELATÓRIO A3", size=11, bold=True, color=NAVY,
        font=FONT_HEAD, align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide, SLIDE_W - Inches(2.6), header_y + Inches(0.42),
        Inches(2.2), Inches(0.32),
        "CICLO DMAIC · 01→02→03→04→05↻", size=10, bold=True,
        color=NAVY_2, font=FONT_HEAD, align=PP_ALIGN.RIGHT,
    )

    # ===== Faixa de destaque dos autores =====
    authors_y = header_y + header_h + Inches(0.10)
    authors_h = Inches(0.78)
    authors_x = Inches(0.45)
    authors_w = SLIDE_W - Inches(0.90)

    # caixa preenchida com leve tint navy
    auth_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, authors_x, authors_y, authors_w, authors_h,
    )
    auth_box.adjustments[0] = 0.18
    auth_box.line.color.rgb = NAVY
    auth_box.line.width = Pt(1.0)
    auth_box.fill.solid()
    auth_box.fill.fore_color.rgb = AUTHORS_BG

    # tag "AUTORES" lateral em navy
    tag_w = Inches(1.25)
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        authors_x, authors_y, tag_w, authors_h,
    )
    tag.adjustments[0] = 0.18
    tag.line.fill.background()
    tag.fill.solid()
    tag.fill.fore_color.rgb = NAVY
    ttf = tag.text_frame
    ttf.margin_left = ttf.margin_right = Emu(0)
    ttf.margin_top = ttf.margin_bottom = Emu(0)
    ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run()
    tr.text = "AUTORES"
    tr.font.name = FONT_HEAD
    tr.font.size = Pt(14)
    tr.font.bold = True
    tr.font.color.rgb = WHITE

    # texto-orientação dos autores (verbatim do original)
    add_text(
        slide,
        authors_x + tag_w + Inches(0.20),
        authors_y, authors_w - tag_w - Inches(0.30), authors_h,
        "Autores: Nome completo, afiliação, e-mail; Máximo 4 autores; "
        "separados por ponto e vírgula (;). Tamanho 12, regular, justificado.",
        size=13, color=NAVY, font=FONT_BODY,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # ===== Grid DMAIC =====
    margin_x = Inches(0.45)
    gutter = Inches(0.32)
    col_w = (SLIDE_W - margin_x * 2 - gutter) // 2

    grid_top = authors_y + authors_h + Inches(0.18)
    footer_h = Inches(1.20)
    grid_bottom = SLIDE_H - footer_h - Inches(0.10)
    grid_h = grid_bottom - grid_top

    left_x = margin_x
    right_x = margin_x + col_w + gutter

    # Esquerda: Definir + Medir (Medir maior por causa da tabela)
    definir_h = Inches(2.85)
    medir_y = grid_top + definir_h + gutter
    medir_h = grid_h - definir_h - gutter

    # Direita: Analisar + Implementar + Controlar
    analisar_h = Inches(2.40)
    implementar_h = Inches(2.85)
    impl_y = grid_top + analisar_h + gutter
    controlar_y = impl_y + implementar_h + gutter
    controlar_h = grid_h - analisar_h - implementar_h - gutter * 2

    # --- 01 DEFINIR ---
    dmaic_box(
        slide, left_x, grid_top, col_w, definir_h,
        numero=1,
        etapa="Definir",
        header_color=YELLOW,
        body_paragraphs=[
            ("Definir:", {"size": 12, "bold": True, "color": NAVY, "space_after": 3}),
            ("Contexto/escopo: Deve conter o escopo do projeto, seus objetivos, "
             "metas e a equipe responsável pela condução das atividades. Fonte "
             "deve ser Arial, 7.",
             {"size": 11, "color": NAVY_2, "space_after": 6}),
            ("Trabalhos aplicados que propõe melhorias em determinado processo "
             "ou produto podem ser submetidos e apresentados no formato do "
             "Relatório A3.",
             {"size": 11, "color": NAVY_2, "space_after": 6}),
            ("O Relatório A3 é uma ferramenta utilizada para identificar "
             "problemas e propor soluções resumidas. O documento deve ser "
             "enxuto. Utilize somente uma página para todo o A3 incluindo "
             "figuras, gráficos e tabelas, seguindo orientações de formatação "
             "e layout neste documento.",
             {"size": 11, "color": NAVY_2, "space_after": 6}),
            ("O A3 deverá seguir o modelo DMAIC, onde os tópicos Definir, e "
             "Medir devem estar na coluna esquerda do documento, e os tópicos "
             "Analisar, Implementar e Controlar devem estar na coluna direita "
             "do documento. O tamanho de cada tópico pode ser ajustado conforme "
             "a necessidade, porém a lateralidade deve ser respeitada.",
             {"size": 11, "color": NAVY_2}),
        ],
    )

    # --- 02 MEDIR ---
    dmaic_box(
        slide, left_x, medir_y, col_w, medir_h,
        numero=2,
        etapa="Medir",
        header_color=GREEN,
        body_paragraphs=[
            ("Medir:", {"size": 12, "bold": True, "color": NAVY, "space_after": 3}),
            ("Deve conter a descrição de coleta de dados, indicadores que "
             "contribuíram para medida do estado prévio ao projeto A3.",
             {"size": 11, "color": NAVY_2, "space_after": 6}),
            ("Tabelas, Figuras e gráficos: Devem ser numerados em ordem "
             "sequencial em algarismos arábicos (01, 02, 03 etc.). O título "
             "das tabelas, figuras e gráficos devem ser inseridos acima deles "
             "e as fontes devem ser inseridas abaixo, ambos em Arial tamanho "
             "10, centralizados, em negrito. Toda imagem, gráfico, figura, "
             "tabela ou qualquer elemento deve ser legível. Exemplo de tabela, "
             "se aplica também a gráficos e figuras:",
             {"size": 11, "color": NAVY_2}),
        ],
    )

    # tabela exemplo (verbatim) — alta contraste
    table_label_y = medir_y + medir_h - Inches(2.05)
    add_text(
        slide, left_x + Inches(0.22), table_label_y, col_w - Inches(0.44),
        Inches(0.28),
        "Tabela 1. Relação das siglas",
        size=11, bold=True, color=NAVY, font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
    )
    tbl_y = table_label_y + Inches(0.34)
    tbl_w = col_w - Inches(0.44)
    tbl_h = Inches(1.05)
    rows, cols = 3, 2
    table_shape = slide.shapes.add_table(
        rows, cols, left_x + Inches(0.22), tbl_y, tbl_w, tbl_h,
    )
    tbl = table_shape.table
    tbl.columns[0].width = Inches(1.5)
    tbl.columns[1].width = tbl_w - Inches(1.5)

    data = [
        ("Siglas", "Órgão"),
        ("Aneel", "Agência Nacional de Energia Elétrica"),
        ("Anvisa", "Agência Nacional de Vigilância Sanitária"),
    ]
    ZEBRA = RGBColor(0xE2, 0xE8, 0xF2)
    for i, (c0, c1) in enumerate(data):
        for j, txt in enumerate((c0, c1)):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = NAVY
                col = WHITE
                bold = True
                align = PP_ALIGN.CENTER
            else:
                cell.fill.fore_color.rgb = ZEBRA if i % 2 == 1 else WHITE
                col = NAVY
                bold = False
                align = PP_ALIGN.LEFT
            tf = cell.text_frame
            tf.margin_left = Inches(0.10)
            tf.margin_right = Inches(0.10)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = txt
            r.font.name = FONT_HEAD if i == 0 else FONT_BODY
            r.font.size = Pt(11)
            r.font.bold = bold
            r.font.color.rgb = col

    add_text(
        slide, left_x + Inches(0.22), tbl_y + tbl_h + Inches(0.06),
        col_w - Inches(0.44), Inches(0.24),
        "Fonte: Autor, ano.",
        size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER,
    )

    # --- 03 ANALISAR (LARANJA) ---
    dmaic_box(
        slide, right_x, grid_top, col_w, analisar_h,
        numero=3,
        etapa="Analisar",
        header_color=ORANGE,
        body_paragraphs=[
            ("Analisar:", {"size": 12, "bold": True, "color": NAVY, "space_after": 3}),
            ("Deve conter a metodologia e as ferramentas utilizadas para "
             "entendimento da causa raiz do estado prévio ao projeto A3.",
             {"size": 11, "color": NAVY_2}),
        ],
    )

    # --- 04 IMPLEMENTAR ---
    dmaic_box(
        slide, right_x, impl_y, col_w, implementar_h,
        numero=4,
        etapa="Implementar",
        header_color=BLUE,
        body_paragraphs=[
            ("Implementar:", {"size": 12, "bold": True, "color": NAVY, "space_after": 3}),
            ("Deve conter as ações utilizadas para eliminar a raiz do problema, "
             "possibilitando a melhoria, o plano de implementação de forma "
             "resumida.",
             {"size": 11, "color": NAVY_2, "space_after": 6}),
            ("Plano de ação / Principais entregas (o que? Quem? Quando? Status "
             "da ação):",
             {"size": 11, "bold": True, "color": NAVY}),
        ],
    )

    # --- 05 CONTROLAR (NAVY — barra escura) ---
    dmaic_box(
        slide, right_x, controlar_y, col_w, controlar_h,
        numero=5,
        etapa="Controlar",
        header_color=NAVY_2,
        is_dark_header=True,
        body_paragraphs=[
            ("Controlar:", {"size": 12, "bold": True, "color": NAVY, "space_after": 3}),
            ("Se houver referências, elas também devem ser colocadas aqui.",
             {"size": 11, "color": NAVY_2}),
        ],
    )

    # ===== Setas indicando o ciclo DMAIC =====
    arrow_w = Inches(0.32)
    arrow_h = Inches(0.32)

    # 01 → 02
    cycle_arrow(
        slide,
        left_x + col_w // 2 - arrow_w // 2,
        grid_top + definir_h + (gutter - arrow_h) // 2,
        arrow_w, arrow_h, direction="down", color=YELLOW,
    )
    # 02 → 03 (gutter horizontal, na altura do header de Analisar)
    cycle_arrow(
        slide,
        margin_x + col_w + (gutter - arrow_w) // 2,
        grid_top + Inches(0.15),
        arrow_w, arrow_h, direction="right", color=GREEN,
    )
    # 03 → 04
    cycle_arrow(
        slide,
        right_x + col_w // 2 - arrow_w // 2,
        grid_top + analisar_h + (gutter - arrow_h) // 2,
        arrow_w, arrow_h, direction="down", color=ORANGE,
    )
    # 04 → 05
    cycle_arrow(
        slide,
        right_x + col_w // 2 - arrow_w // 2,
        impl_y + implementar_h + (gutter - arrow_h) // 2,
        arrow_w, arrow_h, direction="down", color=BLUE,
    )
    # 05 → 01 (retorno)
    return_y = controlar_y + controlar_h - Inches(0.34)
    cycle_arrow(
        slide,
        margin_x + col_w + (gutter - arrow_w) // 2,
        return_y, arrow_w, arrow_h, direction="left", color=NAVY_2,
    )
    add_text(
        slide,
        margin_x, return_y - Inches(0.04),
        col_w + gutter + arrow_w, Inches(0.4),
        "↻ retoma o ciclo",
        size=10, italic=True, bold=True, color=NAVY_2,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # ===== Rodapé =====
    footer_y = SLIDE_H - footer_h
    realizadoras_strip(slide, SLIDE_W, footer_y)

    add_text(
        slide, Inches(0.5), SLIDE_H - Inches(0.32),
        Inches(6), Inches(0.22),
        "encontrogeig.org · 08–10 de julho de 2026 · Rio de Janeiro / Niterói",
        size=8, color=MUTED, align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, SLIDE_W - Inches(6.5), SLIDE_H - Inches(0.32),
        Inches(6), Inches(0.22),
        "4º Seminário em Sistemas de Engenharia de Produção · PPGEP - LabDGE/UFF",
        size=8, color=MUTED, align=PP_ALIGN.RIGHT,
    )

    # faixa decorativa rodapé
    add_color_stripe(slide, SLIDE_W, SLIDE_H - Emu(76200), height=Emu(76200))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK {OUT.name}: {OUT.stat().st_size // 1024} KB")


if __name__ == "__main__":
    build()
