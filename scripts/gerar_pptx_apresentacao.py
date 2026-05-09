"""
Gera o template de apresentação (PPTX) para o vídeo de submissão (Fase 1).

Saída: Modelos/Apresentacao_Video_Encontro_GEI.pptx — 16:9, 8 slides, paleta do site.

Paleta extraída de index.html:
  navy #091136 / #0E1A4A · amarelo #F5C842 · verde #7AC74F · azul #4DA8E0 · ciano #5BC6E5
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
OUT = ROOT / "Modelos" / "Apresentacao_Video_Encontro_GEI.pptx"

# Fundo claro — tipografia em navy, faixas e destaques mantêm a paleta da marca.
BG = RGBColor(0xFF, 0xFF, 0xFF)              # fundo principal
BG_SOFT = RGBColor(0xF5, 0xF7, 0xFB)         # placeholders / cards
BG_BORDER = RGBColor(0xD8, 0xDF, 0xEC)       # bordas suaves
NAVY = RGBColor(0x09, 0x11, 0x36)            # texto forte
NAVY_2 = RGBColor(0x29, 0x33, 0x5C)          # subtítulos
YELLOW = RGBColor(0xC9, 0x9A, 0x14)          # variante de "Governança" legível em branco
YELLOW_BAR = RGBColor(0xF5, 0xC8, 0x42)      # nas faixas
GREEN = RGBColor(0x4F, 0x9C, 0x2F)           # "Estratégia" legível em branco
GREEN_BAR = RGBColor(0x7A, 0xC7, 0x4F)
BLUE = RGBColor(0x2F, 0x7E, 0xB8)            # "Inovação" legível em branco
BLUE_BAR = RGBColor(0x4D, 0xA8, 0xE0)
CYAN = RGBColor(0x2A, 0x9C, 0xC4)
CYAN_BAR = RGBColor(0x5B, 0xC6, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5C, 0x67, 0x81)           # cinza-azulado para corpo

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"

LOGO_EVENTO = ASSETS / "logo-completo.png"
LOGO_EVENTO_LIGHT = ASSETS / "logo-fundo-branco.png"
LOGO_UFF = ASSETS / "logo-uff.png"
LOGO_ABAR = ASSETS / "logos" / "abar.png"
LOGO_LABDGE = ASSETS / "logos" / "labdge.png"
LOGO_GIGS = ASSETS / "logos" / "gigs-unicamp.jpg"


def set_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=WHITE,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_color_stripe(slide, y, *, height=Emu(76200)):
    """Faixa decorativa nas 4 cores do logo (amarelo/verde/azul/ciano)."""
    seg_w = SLIDE_W // 4
    for i, c in enumerate((YELLOW_BAR, GREEN_BAR, BLUE_BAR, CYAN_BAR)):
        seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, seg_w * i, y, seg_w, height)
        seg.line.fill.background()
        seg.fill.solid()
        seg.fill.fore_color.rgb = c


def add_image_safe(slide, path, x, y, *, w=None, h=None):
    if path.exists():
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    return None


def add_uff_badge(slide, x, y, w, h):
    """UFF não tem PNG nos assets — represento como badge tipográfico."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.25)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "UFF"
    run.font.name = FONT_HEAD
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = NAVY


def add_realizadoras_strip(slide, y_top):
    """Tira com os 4 logos institucionais (UFF | ABAR | PPGEP-LabDGE | GIGS)."""
    label = add_text(
        slide, Inches(0.5), y_top, Inches(12.33), Inches(0.28),
        "REALIZAÇÃO", size=10, bold=True,
        color=NAVY_2, align=PP_ALIGN.CENTER,
    )
    label.text_frame.paragraphs[0].runs[0].font.name = FONT_HEAD

    logo_h = Inches(0.55)
    y_logos = y_top + Inches(0.34)
    spacing = Inches(0.5)
    slot_w = Inches(1.7)

    total_w = slot_w * 4 + spacing * 3
    start_x = (SLIDE_W - total_w) // 2

    if LOGO_UFF.exists():
        add_image_safe(slide, LOGO_UFF, start_x, y_logos, h=logo_h)
    else:
        add_uff_badge(slide, start_x, y_logos, slot_w, logo_h)

    add_image_safe(slide, LOGO_ABAR, start_x + (slot_w + spacing) * 1, y_logos, h=logo_h)
    add_image_safe(slide, LOGO_LABDGE, start_x + (slot_w + spacing) * 2, y_logos, h=logo_h)
    add_image_safe(slide, LOGO_GIGS, start_x + (slot_w + spacing) * 3, y_logos, h=logo_h)


def add_footer(slide, page_num, total):
    """Rodapé com mini-logo, paginador e faixa decorativa."""
    add_color_stripe(slide, SLIDE_H - Emu(76200), height=Emu(76200))

    add_image_safe(
        slide, LOGO_EVENTO,
        Inches(0.4), SLIDE_H - Inches(0.62),
        h=Inches(0.32),
    )

    add_text(
        slide, Inches(11.5), SLIDE_H - Inches(0.62), Inches(1.3), Inches(0.3),
        f"{page_num} / {total}", size=10, color=MUTED,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )


def section_slide(prs, *, num, total, titulo, cor_titulo, instrucao,
                  extra_box=False):
    """Slide-padrão de seção: título colorido + área de placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    add_text(
        slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.4),
        f"{num:02d} · {titulo.upper()}",
        size=14, bold=True, color=cor_titulo, font=FONT_HEAD,
    )
    add_text(
        slide, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
        titulo, size=44, bold=True, color=NAVY, font=FONT_HEAD,
    )

    seg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.05), Inches(0.6), Emu(45720)
    )
    seg.line.fill.background()
    seg.fill.solid()
    seg.fill.fore_color.rgb = cor_titulo

    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(2.4),
        Inches(11.93), Inches(4.0) if not extra_box else Inches(2.6),
    )
    placeholder.line.color.rgb = BG_BORDER
    placeholder.line.width = Pt(0.75)
    placeholder.line.dash_style = 7  # tracejado
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = BG_SOFT

    tf = placeholder.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.25)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = instrucao
    run.font.name = FONT_BODY
    run.font.size = Pt(15)
    run.font.color.rgb = MUTED
    run.font.italic = True

    if extra_box:
        figbox = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(5.15),
            Inches(11.93), Inches(1.85),
        )
        figbox.line.color.rgb = BG_BORDER
        figbox.line.width = Pt(0.75)
        figbox.line.dash_style = 7
        figbox.fill.solid()
        figbox.fill.fore_color.rgb = WHITE
        ftf = figbox.text_frame
        ftf.margin_left = Inches(0.3)
        ftf.vertical_anchor = MSO_ANCHOR.MIDDLE
        fp = ftf.paragraphs[0]
        fp.alignment = PP_ALIGN.CENTER
        frun = fp.add_run()
        frun.text = "[ Espaço para figura, gráfico ou tabela — legenda em Calibri 10 ]"
        frun.font.name = FONT_BODY
        frun.font.size = Pt(13)
        frun.font.color.rgb = MUTED
        frun.font.italic = True

    add_footer(slide, num, total)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    global SLIDE_W, SLIDE_H
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    total = 7

    # ===== Slide 1 — Capa =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)

    add_image_safe(
        s, LOGO_EVENTO,
        Inches(4.67), Inches(0.6),
        w=Inches(4.0),
    )

    add_text(
        s, Inches(0.7), Inches(2.7), Inches(11.93), Inches(0.5),
        "FASE 1 · VÍDEO DE APRESENTAÇÃO (até 10 minutos)",
        size=13, bold=True, color=YELLOW, font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.7), Inches(3.2), Inches(11.93), Inches(1.4),
        "[ TÍTULO DO TRABALHO ]",
        size=36, bold=True, color=NAVY, font=FONT_HEAD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        s, Inches(0.7), Inches(4.7), Inches(11.93), Inches(0.5),
        "Autor 1 · Autor 2 · Autor 3 · Autor 4",
        size=18, color=MUTED, font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.7), Inches(5.15), Inches(11.93), Inches(0.4),
        "Instituições · Eixo temático",
        size=14, color=MUTED, font=FONT_BODY,
        align=PP_ALIGN.CENTER,
    )

    add_realizadoras_strip(s, Inches(6.05))
    add_color_stripe(s, SLIDE_H - Emu(76200), height=Emu(76200))

    # ===== Slides 2 a 6 — seções científicas =====
    # Tempos calibrados para totalizar 5–10 min de apresentação (mín. 5'00" · máx. 9'45").
    secoes = [
        ("Introdução", YELLOW,
         "[ Contextualize a pesquisa: problema, relevância prática "
         "(governança / regulação / inovação) e a lacuna que o trabalho endereça. "
         "Use no máximo 3 bullets curtos. Falando entre 45 segundos e 1min30s. ]"),
        ("Objetivo", GREEN,
         "[ Enuncie em 1 frase o objetivo geral. Liste, opcionalmente, "
         "2–3 objetivos específicos. Falando entre 30 e 45 segundos. ]"),
        ("Metodologia", BLUE,
         "[ Descreva a abordagem (qualitativa / quantitativa / mista), "
         "o desenho de pesquisa, fonte de dados, instrumentos e período. "
         "Cite normas / frameworks utilizados. Falando entre 1 e 2 minutos. ]"),
    ]
    for i, (tit, cor, instr) in enumerate(secoes, start=2):
        section_slide(prs, num=i, total=total, titulo=tit, cor_titulo=cor, instrucao=instr)

    section_slide(
        prs, num=5, total=total, titulo="Resultados e discussão", cor_titulo=BLUE,
        instrucao=("[ Apresente os principais resultados e discuta-os à luz da literatura. "
                   "Inclua 1 figura ou tabela na área abaixo. "
                   "Falando entre 2 e 4 minutos. ]"),
        extra_box=True,
    )

    section_slide(
        prs, num=6, total=total, titulo="Considerações finais", cor_titulo=CYAN,
        instrucao=("[ Sintetize as conclusões, limitações da pesquisa e implicações "
                   "para a prática regulatória / inovação governamental. Indique "
                   "trabalhos futuros. Falando entre 45 segundos e 1min30s. ]"),
    )

    # ===== Slide 7 — Encerramento =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)

    add_image_safe(s, LOGO_EVENTO, Inches(4.67), Inches(0.6), w=Inches(4.0))

    add_text(s, Inches(0.7), Inches(2.6), Inches(11.93), Inches(1.0),
             "Obrigado(a)!", size=54, bold=True, color=NAVY, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(3.6), Inches(11.93), Inches(0.5),
             "Perguntas e contato", size=18, color=YELLOW, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.1), Inches(11.93), Inches(0.4),
             "autor@instituicao.br", size=15, color=MUTED, font=FONT_BODY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.5), Inches(11.93), Inches(0.4),
             "encontrogeig.org · 08–10 de julho de 2026 · Rio de Janeiro / Niterói",
             size=14, color=MUTED, font=FONT_BODY, align=PP_ALIGN.CENTER)

    add_realizadoras_strip(s, Inches(5.85))
    add_color_stripe(s, SLIDE_H - Emu(76200), height=Emu(76200))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK {OUT.name}: {OUT.stat().st_size // 1024} KB · {total} slides")


if __name__ == "__main__":
    build()
